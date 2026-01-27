#!/usr/bin/env python3
from __future__ import annotations

import argparse
import json
import sys
from pathlib import Path

try:
    from pptx import Presentation
    from pptx.util import Inches
except ImportError:  # pragma: no cover - runtime guard for missing dependency
    print("Missing dependency: python-pptx. Install requirements.txt first.")
    sys.exit(1)


PROJECT_ROOT = Path(__file__).resolve().parents[1]
FIXTURES_DIR = PROJECT_ROOT / "assets" / "fixtures"
DECK_CATALOG_PATH = FIXTURES_DIR / "clarksoft_deck_catalog.json"
DEFAULT_DECK_PATH = FIXTURES_DIR / "clarksoft_slide_deck.json"
ASSETS_DIR = PROJECT_ROOT / "assets"


def _load_json(path: Path) -> dict:
    try:
        return json.loads(path.read_text(encoding="utf-8"))
    except FileNotFoundError:
        raise FileNotFoundError(f"Missing JSON file: {path}") from None
    except json.JSONDecodeError as exc:
        raise ValueError(f"Invalid JSON in {path}: {exc}") from exc


def _resolve_deck_path(deck_id: str | None) -> Path:
    if not DECK_CATALOG_PATH.exists():
        if DEFAULT_DECK_PATH.exists():
            return DEFAULT_DECK_PATH
        raise FileNotFoundError("Deck catalog and default deck file are missing")

    catalog = _load_json(DECK_CATALOG_PATH)
    resolved_deck_id = deck_id or catalog.get("default_deck_id")
    for deck in catalog.get("decks", []):
        if deck.get("id") == resolved_deck_id and deck.get("file"):
            candidate = FIXTURES_DIR / deck["file"]
            if candidate.exists():
                return candidate

    if DEFAULT_DECK_PATH.exists():
        return DEFAULT_DECK_PATH
    raise FileNotFoundError(f"Deck file not found for deck id {resolved_deck_id!r}")


def _resolve_output_path(raw_output: str | None) -> Path:
    if not raw_output:
        raw_output = "clarksoft_hmi_presenter.pptx"
    output_path = Path(raw_output)
    if output_path.is_absolute():
        output_path.parent.mkdir(parents=True, exist_ok=True)
        return output_path

    exports_dir = PROJECT_ROOT / "assets" / "exports"
    exports_dir.mkdir(parents=True, exist_ok=True)
    return exports_dir / output_path


def _resolve_media_path(media_value: str | None) -> Path | None:
    if not media_value:
        return None
    candidate = Path(media_value)
    if media_value.startswith("/assets/"):
        candidate = ASSETS_DIR / media_value.removeprefix("/assets/")
    elif not candidate.is_absolute():
        candidate = ASSETS_DIR / media_value
    return candidate if candidate.exists() else None


def _add_media(slide, presentation: Presentation, media_path: Path) -> None:
    slide_width = presentation.slide_width
    left = Inches(1)
    top = Inches(2.2)
    width = slide_width - Inches(2)
    try:
        slide.shapes.add_picture(str(media_path), left, top, width=width)
    except Exception:
        # Best-effort: if the image fails, keep the slide text only.
        return


def build_presentation(title: str, subtitle: str, deck: dict) -> Presentation:
    presentation = Presentation()

    title_slide_layout = presentation.slide_layouts[0]
    title_slide = presentation.slides.add_slide(title_slide_layout)
    title_slide.shapes.title.text = title
    title_slide.placeholders[1].text = subtitle

    content_layout = presentation.slide_layouts[1]
    for slide_data in deck.get("slides", []):
        slide = presentation.slides.add_slide(content_layout)
        slide.shapes.title.text = slide_data.get("title", "")

        body_placeholder = slide.shapes.placeholders[1]
        text_frame = body_placeholder.text_frame
        text_frame.text = slide_data.get("body", "") or slide_data.get("notes", "")

        media_path = _resolve_media_path(slide_data.get("media"))
        if media_path:
            _add_media(slide, presentation, media_path)

    return presentation


def main() -> int:
    parser = argparse.ArgumentParser(description="Export a ClarkSoft PPTX deck.")
    parser.add_argument("--output", help="Output file path (.pptx)")
    parser.add_argument("--title", default="ClarkSoft HMI Presenter", help="Title slide heading")
    parser.add_argument("--subtitle", default="ClarkSoft demo deck", help="Title slide subtitle")
    parser.add_argument("--deck", help="Deck id from clarksoft_deck_catalog.json")
    args = parser.parse_args()

    deck_path = _resolve_deck_path(args.deck)
    deck = _load_json(deck_path)
    output_path = _resolve_output_path(args.output)

    presentation = build_presentation(args.title, args.subtitle, deck)
    presentation.save(output_path)

    print(f"Saved PPTX to {output_path}")
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
